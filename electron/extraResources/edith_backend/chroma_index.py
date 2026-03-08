#!/usr/bin/env python3
import os
import re
import csv
import json
import hashlib
import signal
import struct
import zipfile
import tempfile
import time as _time
from pathlib import Path

try:
    from indexing_improvements import (
        IndexCheckpoint, IndexProgressBroadcaster, IngestValidator,
        DuplicateDetector, content_aware_chunk_params, detect_language,
        get_embed_version_info, check_embedding_migration_needed,
        parallel_gemini_embed, enhanced_ocr_extract,
        extract_table_from_text, tables_to_searchable_text,
    )
    _IMPROVEMENTS_AVAILABLE = True
except ImportError:
    _IMPROVEMENTS_AVAILABLE = False

from dotenv import load_dotenv

try:
    # Enable offline model loading from cache when network is unavailable
    import os as _os
    if not _os.environ.get("HF_HUB_OFFLINE"):
        _os.environ.setdefault("HF_HUB_OFFLINE", "0")
    import chromadb  # type: ignore
    from sentence_transformers import SentenceTransformer  # type: ignore
except Exception as e:
    raise SystemExit(f"Missing dependencies for Chroma indexing: {e}")

try:
    from pypdf import PdfReader  # type: ignore
except Exception:
    PdfReader = None

try:
    from docx import Document  # type: ignore
except Exception:
    Document = None

try:
    from pdf2image import convert_from_path  # type: ignore
except Exception:
    convert_from_path = None

try:
    import pytesseract  # type: ignore
except Exception:
    pytesseract = None

try:
    import openpyxl  # type: ignore
except Exception:
    openpyxl = None

try:
    import pandas as pd  # type: ignore
except Exception:
    pd = None

try:
    from dbfread import DBF  # type: ignore
except Exception:
    DBF = None

try:
    import xlrd  # type: ignore
except Exception:
    xlrd = None

try:
    from striprtf.striprtf import rtf_to_text  # type: ignore
except Exception:
    rtf_to_text = None


# ---------- Gemini Embedding API ----------
def _gemini_embed_batch(texts: list, model: str = "gemini-embedding-001", task_type: str = "RETRIEVAL_DOCUMENT") -> list:
    """Embed texts using Google Gemini API. Returns list of vectors or empty list on failure."""
    api_key = os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY") or ""
    if not api_key:
        return []
    try:
        import requests  # type: ignore
    except ImportError:
        return []
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:batchEmbedContents?key={api_key}"
    all_vectors = []
    batch_size = int(os.environ.get("EDITH_GEMINI_BATCH_SIZE", "20"))
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        requests_body = [{"model": f"models/{model}", "content": {"parts": [{"text": t[:8192]}]}, "taskType": task_type} for t in batch]
        success = False
        for retry in range(4):
            try:
                resp = requests.post(url, json={"requests": requests_body}, timeout=120)
                if resp.status_code == 429:
                    wait = 60 * (retry + 1)
                    print(f"  [gemini-embed] 429 rate limit, waiting {wait}s (retry {retry+1}/4)...")
                    _time.sleep(wait)
                    continue
                if resp.status_code != 200:
                    print(f"  [gemini-embed] API error {resp.status_code}: {resp.text[:200]}")
                    return []
                data = resp.json()
                embeddings = data.get("embeddings", [])
                for emb in embeddings:
                    all_vectors.append(emb.get("values", []))
                success = True
                break
            except Exception as e:
                print(f"  [gemini-embed] request failed: {e}")
                return []
        if not success:
            print(f"  [gemini-embed] sub-batch {i//batch_size} failed after 4 retries")
            return []
        if i + batch_size < len(texts):
            _time.sleep(1)
    return all_vectors


def _gemini_embed_queries(texts: list, model: str = "gemini-embedding-001") -> list:
    """Embed query texts using Gemini (uses RETRIEVAL_QUERY task type)."""
    return _gemini_embed_batch(texts, model=model, task_type="RETRIEVAL_QUERY")


# ---------- Hash Manifest ----------
def _load_hash_manifest(manifest_path: Path) -> set:
    """Load known hashes from manifest file (< 1 second vs 60 min from ChromaDB)."""
    if not manifest_path.exists():
        return set()
    try:
        data = json.loads(manifest_path.read_text())
        hashes = set(data.get("hashes", []))
        print(f"  [manifest] loaded {len(hashes)} hashes from {manifest_path.name}")
        return hashes
    except Exception as e:
        print(f"  [manifest] failed to load: {e}")
        return set()


def _save_hash_manifest(manifest_path: Path, hashes: set, docs_root: str, chroma_dir: str):
    """Save known hashes to manifest file for instant startup next time.
    §8.6: Uses atomic write (temp file + rename) for crash safety."""
    try:
        data = {
            "hashes": sorted(hashes),
            "count": len(hashes),
            "docs_root": str(docs_root),
            "chroma_dir": str(chroma_dir),
            "saved_at": _time.strftime("%Y-%m-%dT%H:%M:%S"),
        }
        # §8.6: Atomic write — write to temp, then rename
        tmp_path = manifest_path.with_suffix(".tmp")
        tmp_path.write_text(json.dumps(data, indent=2))
        tmp_path.replace(manifest_path)  # atomic on POSIX
        print(f"  [manifest] saved {len(hashes)} hashes to {manifest_path.name}")
    except Exception as e:
        print(f"  [manifest] failed to save: {e}")


def load_env():
    app_home = Path.home() / "Library" / "Application Support" / "Edith"
    candidates = []
    override = os.environ.get("EDITH_DOTENV_PATH")
    if override:
        candidates.append(Path(override).expanduser())
    candidates.extend([Path(__file__).parent / ".env", Path.cwd() / ".env", app_home / ".env"])

    seen = set()
    for p in candidates:
        key = str(p)
        if key in seen:
            continue
        seen.add(key)
        if p.exists():
            load_dotenv(dotenv_path=p, override=False)


def clean_text(s: str):
    if not s:
        return ""
    s = s.replace("\x00", " ")
    s = re.sub(r"\s+", " ", s)
    return s.strip()


def sha256_file(path: Path):
    h = hashlib.sha256()
    try:
        with path.open("rb") as f:
            for chunk in iter(lambda: f.read(1024 * 1024), b""):
                h.update(chunk)
    except (TimeoutError, OSError, IOError) as e:
        print(f"  [WARN] Cannot hash {path.name}: {e}")
        return None
    return h.hexdigest()


def infer_tag(filename: str):
    stem = Path(filename).stem
    m = re.search(r"#([A-Za-z0-9_-]+)", stem)
    if m:
        return m.group(1)
    m = re.search(r"\[([A-Za-z0-9_-]+)\]", stem)
    if m:
        return m.group(1)
    return ""


def infer_project(rel_path: str):
    parts = Path(rel_path).parts
    if len(parts) >= 2:
        return parts[0]
    return ""


def infer_academic_topic(rel_path: str):
    parts = Path(rel_path).parts
    # If path is canon/Time Series/file.pdf, topic is Time Series
    if len(parts) >= 3:
        return parts[1]
    # If path is inbox/file.pdf, topic is inbox (fallback)
    if len(parts) >= 2:
        return parts[0]
    return "general"


def infer_doc_type(rel_path: str, suffix: str):
    lp = str(rel_path or "").lower()
    ext = str(suffix or "").lower()
    if "chapter" in lp or re.search(r"\bch[0-9]+\b", lp):
        return "thesis_chapter"
    if "slide" in lp or "deck" in lp:
        return "slide"
    if "log" in lp:
        return "log"
    if ext in {".pdf"}:
        return "paper"
    if ext in {".md", ".txt", ".docx", ".doc"}:
        return "note"
    if ext in {".csv", ".tsv", ".xlsx", ".xls"}:
        return "data_table"
    if ext in {".py", ".r", ".sql", ".js", ".ts"}:
        return "code"
    return "document"


def infer_version_stage(filename: str):
    n = clean_text(filename or "").lower()
    if re.search(r"\b(final|camera[-\s]?ready|accepted|published)\b", n):
        return "final"
    if re.search(r"\b(preprint|submission)\b", n):
        return "preprint"
    if re.search(r"\b(draft|wip|working)\b", n):
        return "draft"
    if re.search(r"\b(v[0-9]+|rev[0-9]+|revision)\b", n):
        return "revision"
    return "unknown"


# ── §AI CLASSIFICATION: Gemini-powered auto-tagging ──
_ai_classify_cache: dict = {}  # sha256 -> classification result
_ai_classify_cache_path: Path | None = None


def _load_ai_classify_cache(app_state: Path) -> dict:
    """Load the AI classification cache from disk."""
    global _ai_classify_cache, _ai_classify_cache_path
    _ai_classify_cache_path = app_state / "edith_ai_classify_cache.json"
    if _ai_classify_cache_path.exists():
        try:
            _ai_classify_cache = json.loads(_ai_classify_cache_path.read_text(encoding="utf-8"))
            print(f"  [ai-classify] Loaded cache: {len(_ai_classify_cache)} papers")
        except Exception:
            _ai_classify_cache = {}
    return _ai_classify_cache


def _save_ai_classify_cache():
    """Save the AI classification cache to disk."""
    if _ai_classify_cache_path:
        try:
            _ai_classify_cache_path.write_text(
                json.dumps(_ai_classify_cache, indent=1), encoding="utf-8"
            )
        except Exception:
            pass


def ai_classify_paper(text_sample: str, sha256: str, filename: str = "",
                      existing_topic: str = "") -> dict:
    """Use Gemini to auto-classify a paper's topic, method, country, and theory.

    Args:
        text_sample: First ~2000 chars of extracted text
        sha256: File hash for caching
        filename: Original filename (for context)
        existing_topic: Folder-based topic (may be overridden)

    Returns:
        dict with keys: topic, method, country, theory, doc_type_refined
        Empty strings if classification fails or is unavailable.
    """
    # Check cache first
    if sha256 in _ai_classify_cache:
        return _ai_classify_cache[sha256]

    # Skip if text is too short to classify
    if len(text_sample.strip()) < 200:
        return {"topic": "", "method": "", "country": "", "theory": "", "doc_type_refined": ""}

    # Only classify papers that landed in inbox or have vague topics
    vague_topics = {"inbox", "general", "", "uncategorized", "root"}
    needs_ai = existing_topic.lower() in vague_topics

    # Even if folder-based topic exists, still try to classify method/country
    try:
        import google.generativeai as genai
        prompt = f"""Classify this academic paper. Read the text and return ONLY a JSON object with these fields:
- "topic": The primary academic topic/subfield (e.g., "Survey Design", "Democratic Theory", "International Relations", "Causal Inference", "Political Behavior", "Comparative Politics", "American Politics", "Public Policy")
- "method": The primary research method used (e.g., "Conjoint", "RCT", "DiD", "RDD", "Survey", "Qualitative", "Matching", "IV", "Regression", "Case Study", "Process Tracing", "Content Analysis", "Meta-Analysis", or "" if not applicable)
- "country": The primary country/region studied (e.g., "United States", "Poland", "Mexico", "Europe", "Latin America", or "" if not country-specific)
- "theory": The main theoretical framework (e.g., "Rational Choice", "Institutional Theory", "Democratic Theory", "Behavioralism", or "" if not theory-focused)
- "doc_type_refined": More specific document type (e.g., "empirical_paper", "theoretical_paper", "methods_paper", "review_article", "book_chapter", "working_paper", "dissertation", "policy_brief")

Filename: {filename}
Text excerpt:
{text_sample[:2000]}

Return ONLY valid JSON, no markdown formatting."""

        model = genai.GenerativeModel("gemini-2.0-flash")
        response = model.generate_content(prompt)
        result_text = response.text.strip()

        # Parse JSON response
        if result_text.startswith("```"):
            result_text = result_text.split("```")[1]
            if result_text.startswith("json"):
                result_text = result_text[4:]
        result = json.loads(result_text.strip())

        # Validate and clean
        classification = {
            "topic": str(result.get("topic", "")).strip(),
            "method": str(result.get("method", "")).strip(),
            "country": str(result.get("country", "")).strip(),
            "theory": str(result.get("theory", "")).strip(),
            "doc_type_refined": str(result.get("doc_type_refined", "")).strip(),
        }

        # Cache the result
        _ai_classify_cache[sha256] = classification
        # Save cache periodically
        if len(_ai_classify_cache) % 10 == 0:
            _save_ai_classify_cache()

        return classification

    except ImportError:
        # Gemini not available
        return {"topic": "", "method": "", "country": "", "theory": "", "doc_type_refined": ""}
    except Exception as e:
        print(f"  [ai-classify] warning: {e}")
        return {"topic": "", "method": "", "country": "", "theory": "", "doc_type_refined": ""}


def infer_citation(filename: str):
    stem = Path(filename).stem
    cleaned = re.sub(r"[_]+|[-]+", " ", stem)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    year_match = re.search(r"(19|20)\d{2}", cleaned)
    year = year_match.group(0) if year_match else ""
    author = ""
    title = cleaned
    if year_match:
        before, _, after = cleaned.partition(year)
        before = before.strip(" -_()")
        after = after.strip(" -_()")
        if before:
            author = before.split(",")[0].strip()
        if after:
            title = after
    return title, author, year


def normalize_family_component(value: str):
    text = clean_text(value or "").lower()
    text = re.sub(r"[^a-z0-9]+", "_", text).strip("_")
    return text[:80]


def infer_doc_family(title: str, author: str, year: str, rel_path: str):
    parts = []
    if author:
        parts.append(normalize_family_component(author.split()[0]))
    if year and re.fullmatch(r"(19|20)\d{2}", str(year).strip()):
        parts.append(str(year).strip())
    if title:
        parts.append(normalize_family_component(title)[:50])
    family = "_".join([p for p in parts if p])
    if not family:
        family = normalize_family_component(Path(rel_path).stem)
    return family or normalize_family_component(rel_path)

def extract_figure_table_markers(text: str, limit: int = 6):
    """
    Detects markers like 'Figure 1', 'Table 2', 'Tab. 1', etc.
    """
    if len(text) > 100000:
        return []
    markers = []
    # Standard academic markers
    matches = re.finditer(r"(?:Figure|Fig\.|Table|Tab\.|Diagram)\s+(\d+[:\.]?\d*)", text, flags=re.I)
    for m in matches:
        markers.append(m.group(0).strip())
    # Markdown-style tables
    if "|" in text and "-" in text and "+" in text:
        markers.append("MarkdownTable")
    
    # Prune and limit
    out = []
    for m in markers:
        if m.lower() not in {x.lower() for x in out}:
            out.append(m)
        if len(out) >= limit:
            break
    return out


def extract_equation_markers(text: str, limit: int = 6):
    markers = []
    for m in re.finditer(r"\b([A-Za-z][A-Za-z0-9_]{0,20}\s*=\s*[^.;\n]{3,100})", (text or "")[:100000]):
        eq = clean_text(m.group(1))
        if not eq:
            continue
        key = eq.lower()
        if key in {x.lower() for x in markers}:
            continue
        markers.append(eq)
        if len(markers) >= limit:
            break
    return markers


def extract_caption_blocks(text: str, page: int | None):
    """
    Pull likely figure/table captions as separate searchable blocks.
    Keeps recall high for results-style questions.
    """
    out = []
    if not text:
        return out
    for m in re.finditer(r"\b(Figure|Fig\.|Table)\s*([A-Za-z0-9.\-]+)\s*[:.\-]?\s*([^.;]{18,260})", text, flags=re.I):
        label = clean_text(f"{m.group(1)} {m.group(2)}")
        cap = clean_text(m.group(3))
        if not cap:
            continue
        snippet = clean_text(f"{label}: {cap}")
        out.append(
            {
                "page": page,
                "section_heading": "Figure/Table caption",
                "text": snippet,
            }
        )
    # de-duplicate while keeping order
    seen = set()
    deduped = []
    for row in out:
        key = row.get("text", "").lower()
        if not key or key in seen:
            continue
        seen.add(key)
        deduped.append(row)
    return deduped[:8]


def infer_citation_from_pdf(path: Path):
    if PdfReader is None:
        return "", "", "", "filename"
    try:
        reader = PdfReader(str(path))
    except Exception:
        return "", "", "", "filename"

    title = ""
    author = ""
    year = ""
    source = "filename"

    meta = reader.metadata or {}
    title = clean_text(str(meta.get("/Title", "") or meta.get("title", "") or ""))
    author = clean_text(str(meta.get("/Author", "") or meta.get("author", "") or ""))
    date_candidates = [
        str(meta.get("/CreationDate", "") or ""),
        str(meta.get("/ModDate", "") or ""),
    ]
    for raw in date_candidates:
        m = re.search(r"(19|20)\d{2}", raw)
        if m:
            year = m.group(0)
            break

    if title or author or year:
        source = "pdf_metadata"

    first_page_text = ""
    try:
        if reader.pages:
            first_page_text = clean_text(reader.pages[0].extract_text() or "")
    except Exception:
        first_page_text = ""

    if first_page_text:
        if not year:
            m = re.search(r"\b(19|20)\d{2}\b", first_page_text)
            if m:
                year = m.group(0)
        lines = [clean_text(x) for x in first_page_text.split("  ") if clean_text(x)]
        if not title:
            for line in lines[:12]:
                if len(line) < 15:
                    continue
                lower = line.lower()
                if "abstract" in lower or "introduction" in lower:
                    continue
                if re.search(r"https?://", lower):
                    continue
                title = line[:220]
                source = "pdf_first_page"
                break
        if not author:
            author_match = re.search(
                r"\b([A-Z][A-Za-z'`-]+(?:\s+[A-Z][A-Za-z'`-]+){1,3})\b",
                first_page_text,
            )
            if author_match:
                author = author_match.group(1)
                if source == "filename":
                    source = "pdf_first_page"

    return title, author, year, source


def _paragraph_similarity(p1: str, p2: str) -> float:
    """Lightweight cosine similarity between two paragraphs using word overlap.
    No embedding API calls — just fast TF-IDF-like comparison."""
    words1 = set(re.findall(r'\b\w{3,}\b', p1.lower()))
    words2 = set(re.findall(r'\b\w{3,}\b', p2.lower()))
    if not words1 or not words2:
        return 0.0
    intersection = words1 & words2
    union = words1 | words2
    return len(intersection) / len(union) if union else 0.0


def _is_quality_chunk(text: str) -> bool:
    """§5.4: Reject garbage chunks (mostly whitespace, numbers, or formatting)."""
    if not text or len(text.strip()) < 20:
        return False
    stripped = text.strip()
    # Reject if > 80% non-alpha characters
    alpha_chars = sum(1 for c in stripped if c.isalpha())
    if len(stripped) > 0 and alpha_chars / len(stripped) < 0.2:
        return False
    # Reject if fewer than 3 words
    if len(stripped.split()) < 3:
        return False
    return True


def _chunk_quality_score(text: str) -> float:
    """§4.0: Compute chunk information density (0.0-1.0).
    Higher = denser, more useful for retrieval. Used to boost relevant chunks."""
    if not text or len(text.strip()) < 20:
        return 0.0
    stripped = text.strip()
    # Factor 1: Alpha density (penalize chunks that are mostly numbers/formatting)
    alpha_ratio = sum(1 for c in stripped if c.isalpha()) / max(len(stripped), 1)
    # Factor 2: Unique word ratio (penalize repetitive text)
    words = stripped.lower().split()
    unique_ratio = len(set(words)) / max(len(words), 1) if words else 0
    # Factor 3: Sentence count (more sentences = probably more info)
    sentences = len(re.findall(r'[.!?]+', stripped))
    sentence_factor = min(sentences / 5.0, 1.0)
    # Factor 4: Contains citation-like patterns (Author, Year) or [S#]
    has_citations = 1.0 if re.search(r'\([A-Z][a-z]+.*?(19|20)\d{2}\)', stripped) else 0.5
    # Weighted combination
    score = (alpha_ratio * 0.3 + unique_ratio * 0.3 +
             sentence_factor * 0.2 + has_citations * 0.2)
    return round(min(max(score, 0.0), 1.0), 3)


def chunk_text(text: str, chunk_size: int, overlap: int):
    """
    Semantic chunking: splits at paragraph boundaries where topic changes.
    Uses word-overlap similarity between adjacent paragraphs to detect
    natural topic shifts. Chunks stay together when paragraphs discuss
    the same topic and split where the topic changes.

    §5.4: Filters out garbage chunks (mostly whitespace/numbers).
    §5.9: Merges tiny chunks (< 100 chars) into adjacent chunks.
    """
    text = (text or "").strip()
    if not text:
        return []
    
    # 1. Split by paragraphs
    paragraphs = re.split(r'\n\s*\n', text)
    paragraphs = [p.strip() for p in paragraphs if p.strip()]
    
    if not paragraphs:
        return []
    
    # Short text: return as-is
    total_len = sum(len(p) for p in paragraphs)
    if total_len <= chunk_size:
        return ["\n\n".join(paragraphs)]
    
    # 2. Compute similarity between adjacent paragraphs
    similarity_threshold = 0.15  # Below this = topic shift = chunk boundary
    
    # 3. Group paragraphs by semantic similarity
    raw_chunks = []
    current_group = [paragraphs[0]]
    current_len = len(paragraphs[0])
    
    for i in range(1, len(paragraphs)):
        para = paragraphs[i]
        
        # If this single paragraph is huge, we must split it by characters
        if len(para) > chunk_size:
            # Flush current group
            if current_group:
                raw_chunks.append("\n\n".join(current_group))
                current_group = []
                current_len = 0
            # Split the giant paragraph
            j = 0
            while j < len(para):
                end = min(len(para), j + chunk_size)
                raw_chunks.append(para[j:end])
                if end >= len(para):
                    break
                j = max(0, end - overlap)
            continue
        
        # Check semantic similarity with previous paragraph
        sim = _paragraph_similarity(current_group[-1], para)
        
        # Decide whether to create a boundary
        is_topic_shift = sim < similarity_threshold
        would_exceed_size = current_len + len(para) + 2 > chunk_size
        
        if would_exceed_size or (is_topic_shift and current_len > chunk_size // 3):
            # Create chunk boundary
            raw_chunks.append("\n\n".join(current_group))
            # Overlap: keep last paragraph if small enough for context
            if len(current_group[-1]) < overlap and not is_topic_shift:
                current_group = [current_group[-1], para]
                current_len = len(current_group[0]) + len(para) + 2
            else:
                current_group = [para]
                current_len = len(para)
        else:
            # Same topic — keep together
            current_group.append(para)
            current_len += len(para) + 2
            
    if current_group:
        raw_chunks.append("\n\n".join(current_group))

    # §5.9: Merge tiny chunks (< 100 chars) into adjacent chunks
    MIN_CHUNK = 100
    merged = []
    for ch in raw_chunks:
        if merged and len(ch) < MIN_CHUNK:
            merged[-1] = merged[-1] + "\n\n" + ch
        else:
            merged.append(ch)
    # If last chunk became tiny after loop, merge backward
    if len(merged) > 1 and len(merged[-1]) < MIN_CHUNK:
        merged[-2] = merged[-2] + "\n\n" + merged[-1]
        merged.pop()

    # §5.4: Filter out garbage chunks
    return [ch for ch in merged if _is_quality_chunk(ch)]


def flatten_if_nested(values):
    if not values:
        return []
    if isinstance(values, list) and values and isinstance(values[0], list):
        return values[0]
    return values


def load_vault_file_manifest(docs_root: Path):
    override = (os.environ.get("EDITH_VAULT_FILE_MANIFEST") or "").strip()
    path = Path(override).expanduser() if override else (docs_root / "vault_sync" / "_manifests" / "vault_file_manifest.json")
    if not path.exists():
        return {}
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
        return payload if isinstance(payload, dict) else {}
    except Exception:
        return {}


def join_manifest_field(value):
    if isinstance(value, list):
        out = []
        for item in value:
            s = clean_text(str(item or ""))
            if s and s not in out:
                out.append(s)
        return ",".join(out)
    return clean_text(str(value or ""))


def _extract_pdf_fitz(path: Path):
    """Extract PDF pages using PyMuPDF (fitz) — much cleaner than pypdf."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return None
    try:
        doc = fitz.open(str(path))
    except Exception:
        return None
    pages = []
    for idx, page in enumerate(doc, start=1):
        try:
            txt = clean_text(page.get_text("text") or "")
        except Exception:
            txt = ""
        if txt:
            pages.append((idx, txt, False))
    doc.close()
    return pages if pages else None


def extract_pdf_pages(path: Path, ocr_on_ingest: bool = False, ocr_max_pages: int = 20):
    """§4.10: Per-page error recovery — if a page fails, continue with next page."""
    # Try PyMuPDF first (better font handling, no "wrong pointing object" warnings)
    fitz_result = _extract_pdf_fitz(path)
    if fitz_result:
        return fitz_result

    # Fall back to pypdf
    if PdfReader is None:
        return []
    try:
        reader = PdfReader(str(path))
    except Exception:
        return []
    pages = []
    missing_pages = []
    page_errors = 0
    for idx, page in enumerate(reader.pages, start=1):
        try:
            txt = clean_text(page.extract_text() or "")
        except Exception as e:
            # §4.10: Log and continue instead of aborting the whole file
            page_errors += 1
            print(f"  [pdf] page {idx} extraction error ({path.name}): {e}")
            txt = ""
        if txt:
            pages.append((idx, txt, False))
        else:
            missing_pages.append(idx)

    if page_errors > 0:
        print(f"  [pdf] {path.name}: {page_errors} page(s) had extraction errors, {len(pages)} pages OK")

    if (
        missing_pages
        and ocr_on_ingest
        and convert_from_path is not None
        and pytesseract is not None
    ):
        for page_idx in missing_pages[: max(0, int(ocr_max_pages))]:
            try:
                images = convert_from_path(
                    str(path),
                    dpi=200,
                    first_page=int(page_idx),
                    last_page=int(page_idx),
                )
                if not images:
                    continue
                ocr_text = clean_text(pytesseract.image_to_string(images[0]) or "")
                if ocr_text:
                    pages.append((int(page_idx), ocr_text, True))
            except Exception:
                continue
    pages.sort(key=lambda x: int(x[0]))
    return pages


def guess_section_heading(text: str):
    if not text:
        return ""
    raw = text.strip()
    lines = [clean_text(x) for x in raw.split("\n") if clean_text(x)]
    if not lines:
        return ""
    first = lines[0]
    if len(first) > 90:
        return ""
    words = first.split()
    if len(words) < 2 or len(words) > 12:
        return ""
    lowered = first.lower()
    if lowered in {"abstract", "introduction", "references", "bibliography", "appendix"}:
        return first
    upper_ratio = sum(1 for c in first if c.isalpha() and c.isupper()) / float(max(1, sum(1 for c in first if c.isalpha())))
    if upper_ratio > 0.65:
        return first.title()
    if first[0].isupper():
        return first
    return ""


def extract_markdown_blocks(path: Path):
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return []
    lines = raw.splitlines()
    blocks = []
    current_heading = "Document"
    current_lines = []

    def flush():
        text = clean_text("\n".join(current_lines))
        if text:
            blocks.append({"page": None, "section_heading": current_heading, "text": text})

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("#"):
            flush()
            current_lines = []
            heading = re.sub(r"^#+\s*", "", stripped).strip()
            current_heading = heading or "Section"
            continue
        current_lines.append(line)

    flush()
    return blocks


def extract_docx_blocks(path: Path):
    if Document is None:
        return []
    try:
        d = Document(str(path))
    except Exception:
        return []

    blocks = []
    current_heading = "Document"
    current_lines = []

    def flush():
        text = clean_text("\n".join(current_lines))
        if text:
            blocks.append({"page": None, "section_heading": current_heading, "text": text})

    for para in d.paragraphs:
        text = clean_text(para.text or "")
        if not text:
            continue
        style_name = ""
        try:
            style_name = (para.style.name or "").strip().lower()
        except Exception:
            style_name = ""

        if style_name.startswith("heading"):
            flush()
            current_lines = []
            current_heading = text
            continue
        current_lines.append(text)

    flush()
    return blocks


def extract_xlsx_blocks(path: Path):
    """Extract text from Excel .xlsx spreadsheets."""
    if openpyxl is None:
        return []
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return []
    blocks = []
    for sheet_name in wb.sheetnames:
        try:
            ws = wb[sheet_name]
            rows_text = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_count += 1
                if row_count > 500:
                    rows_text.append(f"... ({row_count}+ rows total)")
                    break
                cells = [clean_text(str(c or "")) for c in row]
                line = " | ".join([c for c in cells if c])
                if line:
                    rows_text.append(line)
            text = "\n".join(rows_text)
            if text:
                blocks.append({"page": None, "section_heading": f"Sheet: {sheet_name}", "text": text})
        except Exception:
            continue
    try:
        wb.close()
    except Exception:
        pass
    return blocks


def extract_xls_blocks(path: Path):
    """Extract text from old Excel .xls spreadsheets."""
    if xlrd is None:
        return []
    try:
        wb = xlrd.open_workbook(str(path))
    except Exception:
        return []
    blocks = []
    for sheet in wb.sheets():
        try:
            rows_text = []
            for row_idx in range(min(sheet.nrows, 500)):
                cells = [clean_text(str(sheet.cell_value(row_idx, col) or "")) for col in range(sheet.ncols)]
                line = " | ".join([c for c in cells if c])
                if line:
                    rows_text.append(line)
            if sheet.nrows > 500:
                rows_text.append(f"... ({sheet.nrows} rows total)")
            text = "\n".join(rows_text)
            if text:
                blocks.append({"page": None, "section_heading": f"Sheet: {sheet.name}", "text": text})
        except Exception:
            continue
    return blocks


def extract_dta_blocks(path: Path):
    """Extract text from Stata .dta data files."""
    if pd is None:
        return []
    try:
        reader = pd.io.stata.StataReader(str(path))
        var_labels = {}
        try:
            var_labels = reader.variable_labels()
        except Exception:
            pass
        try:
            df = pd.read_stata(str(path), iterator=False)
        except Exception:
            df = pd.read_stata(str(path), convert_categoricals=False)

        parts = []
        col_info = []
        for col in df.columns:
            label = var_labels.get(col, "")
            col_info.append(f"{col}: {label}" if label else col)
        parts.append("Variables: " + ", ".join(col_info))
        parts.append(f"Observations: {len(df)} rows, {len(df.columns)} columns")

        try:
            desc = df.describe(include='all').to_string()
            parts.append(f"Summary Statistics:\n{desc}")
        except Exception:
            pass

        try:
            sample = df.head(20).to_string()
            parts.append(f"Sample Data (first 20 rows):\n{sample}")
        except Exception:
            pass

        text = "\n\n".join(parts)
        reader.close()
        return [{"page": None, "section_heading": "Stata Dataset", "text": text}] if text else []
    except Exception:
        return []


def extract_dbf_blocks(path: Path):
    """Extract text from dBASE .dbf files (GIS attribute tables)."""
    if DBF is None:
        # Fallback: try reading as binary to get field names
        try:
            with open(path, 'rb') as f:
                header = f.read(32)
                if len(header) < 32:
                    return []
                num_records = struct.unpack('<I', header[4:8])[0]
                header_size = struct.unpack('<H', header[8:10])[0]
                num_fields = (header_size - 33) // 32
                fields = []
                for i in range(min(num_fields, 100)):
                    field_data = f.read(32)
                    if len(field_data) < 32 or field_data[0] == 0x0D:
                        break
                    name = field_data[:11].split(b'\x00')[0].decode('ascii', errors='ignore')
                    if name:
                        fields.append(name)
                text = f"dBASE table: {num_records} records, {len(fields)} fields\nFields: {', '.join(fields)}"
                return [{"page": None, "section_heading": "GIS Attribute Table", "text": text}]
        except Exception:
            return []
    try:
        table = DBF(str(path), encoding='utf-8', ignore_missing_memofile=True)
        fields = [f.name for f in table.fields]
        parts = [f"Fields: {', '.join(fields)}"]
        parts.append(f"Records: {len(table)}")
        sample_rows = []
        for i, record in enumerate(table):
            if i >= 20:
                break
            row = " | ".join([f"{k}={clean_text(str(v or ''))}" for k, v in record.items() if v])
            if row:
                sample_rows.append(row)
        if sample_rows:
            parts.append("Sample:\n" + "\n".join(sample_rows))
        text = "\n\n".join(parts)
        return [{"page": None, "section_heading": "GIS Attribute Table", "text": text}] if text else []
    except Exception:
        return []


def extract_rtf_blocks(path: Path):
    """Extract text from .rtf files."""
    if rtf_to_text is None:
        return []
    try:
        raw = path.read_bytes().decode("utf-8", errors="ignore")
        text = clean_text(rtf_to_text(raw))
        if text:
            return [{"page": None, "section_heading": "Document", "text": text}]
    except Exception:
        pass
    return []


def extract_zip_contents(path: Path, ocr_on_ingest: bool = False, ocr_max_pages: int = 20):
    """Extract and index files within a .zip archive."""
    try:
        if not zipfile.is_zipfile(str(path)):
            return []
        blocks = []
        with zipfile.ZipFile(str(path), 'r') as zf:
            names = zf.namelist()
            blocks.append({"page": None, "section_heading": "Archive Contents",
                          "text": f"ZIP archive with {len(names)} files: " + ", ".join(names[:50])})
            extracted = 0
            for name in names:
                if extracted >= 50:
                    break
                info = zf.getinfo(name)
                if info.file_size > 50 * 1024 * 1024:
                    continue
                if info.is_dir():
                    continue
                ext = Path(name).suffix.lower()
                if ext not in {".pdf", ".txt", ".md", ".csv", ".json", ".docx", ".py", ".r", ".do", ".xml", ".tex"}:
                    continue
                try:
                    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                        tmp.write(zf.read(name))
                        tmp_path = Path(tmp.name)
                    inner_blocks = extract_text_blocks(tmp_path, ocr_on_ingest=ocr_on_ingest, ocr_max_pages=ocr_max_pages)
                    for b in inner_blocks:
                        b["section_heading"] = f"[{name}] {b.get('section_heading', '')}"
                    blocks.extend(inner_blocks)
                    extracted += 1
                except Exception:
                    pass
                finally:
                    try:
                        tmp_path.unlink()
                    except Exception:
                        pass
        return blocks
    except Exception:
        return []


def extract_gis_metadata_stub(path: Path):
    """Create a metadata-only stub for binary GIS, image, and graph files.
    Indexes just filename + path + size so Edith knows the file exists."""
    try:
        size_kb = round(path.stat().st_size / 1024, 1)
        ext = path.suffix.lower()
        type_map = {
            ".shp": "Shapefile geometry", ".shx": "Shapefile index", ".prj": "Projection file",
            ".dbf": "Attribute table", ".cpg": "Code page file",
            ".gdbtable": "Geodatabase table", ".gdbtablx": "Geodatabase index",
            ".gdbindexes": "Geodatabase indexes", ".atx": "Attribute index",
            ".sbn": "Spatial index", ".sbx": "Spatial index",
            ".spx": "Spatial index", ".horizon": "Horizon data",
            ".freelist": "Freelist data", ".cfs": "Compound file",
            ".gph": "Stata graph", ".jpg": "Image (JPEG)", ".png": "Image (PNG)",
            ".gif": "Image (GIF)", ".tif": "Image (TIFF)", ".tiff": "Image (TIFF)",
        }
        file_type = type_map.get(ext, f"Binary file ({ext})")
        # Include parent folder for GIS context (e.g., the .gdb folder name)
        parent_folder = path.parent.name
        text = f"File: {path.name}\nType: {file_type}\nFolder: {parent_folder}\nSize: {size_kb} KB"
        return [{"page": None, "section_heading": file_type, "text": text}]
    except Exception:
        return []


def extract_text_blocks(path: Path, ocr_on_ingest: bool = False, ocr_max_pages: int = 20):
    ext = path.suffix.lower()
    if ext == ".pdf":
        out = []
        for p, t, ocr_used in extract_pdf_pages(path, ocr_on_ingest=ocr_on_ingest, ocr_max_pages=ocr_max_pages):
            out.append({"page": p, "section_heading": guess_section_heading(t), "text": t, "ocr_used": bool(ocr_used)})
            out.extend(extract_caption_blocks(t, p))
        if not out:
            print(f"  [WARN] PDF produced zero text blocks: {path.name}")
        return out
    if ext == ".docx":
        try:
            blocks = extract_docx_blocks(path)
            if not blocks:
                print(f"  [WARN] DOCX produced zero blocks: {path.name}")
            return blocks
        except Exception as e:
            print(f"  [ERROR] DOCX extraction failed for {path.name}: {e}")
            return []
    # Data files → metadata stub only (don't pollute vector store with data rows)
    if ext in {".xlsx", ".xls", ".dta", ".csv", ".tsv", ".sav", ".por"}:
        return extract_gis_metadata_stub(path)
    if ext == ".dbf":
        return extract_gis_metadata_stub(path)  # GIS attribute table — metadata only
    # Jupyter notebooks — extract markdown + code cells
    if ext == ".ipynb":
        try:
            nb = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
            cells = nb.get("cells", [])
            blocks = []
            for idx, cell in enumerate(cells, start=1):
                cell_type = cell.get("cell_type", "")
                source = "".join(cell.get("source", []))
                if not source.strip():
                    continue
                if cell_type == "markdown":
                    heading = source.strip().split("\n")[0][:80]
                    text = clean_text(source)
                elif cell_type == "code":
                    heading = f"Code cell {idx}"
                    # Include outputs if present
                    outputs = []
                    for out in cell.get("outputs", []):
                        if "text" in out:
                            outputs.append("".join(out["text"]))
                        elif "data" in out and "text/plain" in out["data"]:
                            outputs.append("".join(out["data"]["text/plain"]))
                    text = clean_text(source)
                    if outputs:
                        text += "\n\nOutput:\n" + clean_text("\n".join(outputs[:5]))
                else:
                    continue
                if text:
                    blocks.append({"page": idx, "section_heading": heading, "text": text})
            if not blocks:
                print(f"  [WARN] Notebook produced zero blocks: {path.name} ({len(cells)} cells)")
            return blocks
        except Exception as e:
            print(f"  [ERROR] Notebook extraction failed for {path.name}: {e}")
            return []
    # R Markdown — extract prose and code chunks
    if ext == ".rmd":
        try:
            raw = path.read_text(encoding="utf-8", errors="ignore")
            blocks = []
            # Split on R code chunks ```{r ...} ... ```
            parts = re.split(r'```\{[^}]*\}', raw)
            code_chunks = re.findall(r'```\{[^}]*\}([\s\S]*?)```', raw)
            for i, prose in enumerate(parts):
                prose = clean_text(prose)
                if prose:
                    heading = prose.strip().split("\n")[0][:80]
                    blocks.append({"page": i + 1, "section_heading": heading, "text": prose})
                if i < len(code_chunks):
                    code_text = clean_text(code_chunks[i])
                    if code_text:
                        blocks.append({"page": i + 1, "section_heading": f"R code chunk {i + 1}", "text": code_text})
            if not blocks:
                print(f"  [WARN] RMarkdown produced zero blocks: {path.name}")
            return blocks
        except Exception as e:
            print(f"  [ERROR] RMarkdown extraction failed for {path.name}: {e}")
            return []
    if ext == ".rtf":
        return extract_rtf_blocks(path)
    if ext == ".zip":
        return extract_zip_contents(path, ocr_on_ingest=ocr_on_ingest, ocr_max_pages=ocr_max_pages)
    if ext == ".md":
        return extract_markdown_blocks(path)
    # PPTX — slide-level chunks
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            blocks = []
            for idx, slide in enumerate(prs.slides, start=1):
                texts = []
                title = f"Slide {idx}"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                    if shape.has_text_frame and shape == slide.shapes.title:
                        title = shape.text.strip()[:80] or title
                slide_text = clean_text("\n".join(texts))
                if slide_text:
                    blocks.append({"page": idx, "section_heading": title, "text": slide_text})
            if not blocks:
                print(f"  [WARN] PPTX produced zero blocks: {path.name}")
            return blocks
        except Exception as e:
            print(f"  [ERROR] PPTX extraction failed for {path.name}: {e}")
            return []
    # Code/script files → metadata stub only (analysis code, not academic prose)
    if ext in {".r", ".py", ".do", ".sql", ".js", ".ts", ".sas", ".sps", ".ado", ".mata"}:
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
            if not text.strip():
                print(f"  [WARN] Empty script file: {path.name}")
                return []
            blocks = []
            current_block = []
            current_heading = "Preamble"
            for line in text.split("\n"):
                stripped = line.strip()
                is_boundary = False
                if ext == ".py" and (stripped.startswith("def ") or stripped.startswith("class ")):
                    is_boundary = True
                    current_heading = stripped.split("(")[0].replace("def ", "").replace("class ", "")
                elif ext == ".r" and ("<-" in stripped and "function" in stripped):
                    is_boundary = True
                    current_heading = stripped.split("<-")[0].strip()
                elif ext == ".do" and stripped.lower().startswith(("program ", "capture ", "forvalues", "foreach ")):
                    is_boundary = True
                    current_heading = stripped[:60]
                elif ext in {".sas", ".sps"} and stripped.upper().startswith(("PROC ", "DATA ", "RUN;", "COMPUTE ")):
                    is_boundary = True
                    current_heading = stripped[:60]
                elif ext == ".sql" and stripped.upper().startswith(("CREATE ", "SELECT ", "INSERT ", "ALTER ")):
                    is_boundary = True
                    current_heading = stripped[:60]
                if is_boundary and current_block:
                    block_text = clean_text("\n".join(current_block))
                    if block_text:
                        blocks.append({"page": None, "section_heading": current_heading, "text": block_text})
                    current_block = []
                current_block.append(line)
            if current_block:
                block_text = clean_text("\n".join(current_block))
                if block_text:
                    blocks.append({"page": None, "section_heading": current_heading, "text": block_text})
            return blocks if blocks else [{"page": None, "section_heading": "Script", "text": clean_text(text)}]
        except Exception as e:
            print(f"  [ERROR] Script extraction failed for {path.name}: {e}")
            return []
    # GIS/image binary files — metadata stub only
    gis_binary_exts = {
        ".shp", ".shx", ".gdbtable", ".gdbtablx", ".gdbindexes",
        ".atx", ".sbn", ".sbx", ".spx", ".horizon", ".freelist", ".cfs",
        ".gph", ".jpg", ".png", ".gif", ".tif", ".tiff",
    }
    if ext in gis_binary_exts or ext in {".dbf", ".prj", ".cpg"}:
        return extract_gis_metadata_stub(path)
    # SMCL (Stata log format) — strip SMCL tags
    if ext == ".smcl":
        try:
            raw = path.read_text(encoding="utf-8", errors="ignore")
            # Strip SMCL formatting tags like {txt}, {res}, etc.
            text = re.sub(r'\{[^}]*\}', '', raw)
            text = clean_text(text)
            if text:
                return [{"page": None, "section_heading": "Stata Log", "text": text}]
            print(f"  [WARN] SMCL file produced empty text: {path.name}")
        except Exception as e:
            print(f"  [ERROR] SMCL extraction failed for {path.name}: {e}")
            return []
    # All other text-like files
    text_exts = {
        ".txt", ".json", ".jsonl", ".xml", ".html", ".htm",
        ".tex", ".bib", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
        ".sh", ".bash", ".sthlp", ".log", ".note",
    }
    if ext in text_exts:
        try:
            text = clean_text(path.read_text(encoding="utf-8", errors="ignore"))
            if text:
                section = guess_section_heading(text) or "Document"
                return [{"page": None, "section_heading": section, "text": text}]
            print(f"  [WARN] Text file produced empty content: {path.name}")
        except Exception as e:
            print(f"  [ERROR] Text file extraction failed for {path.name}: {e}")
            return []
    print(f"  [WARN] Unsupported extension: {ext} for {path.name}")
    return []


# ──────────────────────────────────────────────────────────────────────
#  Decomposed indexing pipeline: setup → scan_and_index → finalize
# ──────────────────────────────────────────────────────────────────────

# GIS/binary extensions that should get metadata-only stubs (no embedding)
GIS_BINARY_EXTS = {
    ".shp", ".shx", ".gdbtable", ".gdbtablx", ".gdbindexes",
    ".atx", ".sbn", ".sbx", ".spx", ".horizon", ".freelist", ".cfs",
    ".gph", ".jpg", ".png", ".gif", ".tif", ".tiff",
    ".gpkg",  # GeoPackage (binary SQLite)
    # GIS-adjacent — metadata only, do NOT embed (pollutes academic vectors)
    ".dbf", ".prj", ".cpg",
}

VALID_EXT = {
    # Documents
    ".pdf", ".txt", ".md", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".rtf",
    # Data
    ".json", ".jsonl", ".csv", ".tsv", ".dta", ".sav", ".rds",
    # Code
    ".py", ".sql", ".r", ".js", ".ts", ".do", ".sas", ".sps", ".ado", ".mata", ".sthlp",
    # Notebooks
    ".ipynb", ".rmd",
    # Markup & config
    ".xml", ".html", ".htm", ".tex", ".bib", ".yaml", ".yml", ".toml",
    # Logs
    ".log", ".smcl",
    # Archives — REMOVED: .zip was chunking binaries into garbage (179 chunks from a .app zip)
    # ".zip",  # Disabled: produces garbage chunks from binary archives
    # Misc text
    ".note",
    # GIS text formats (indexable as text)
    ".geojson", ".kml", ".gpx",
    # GIS/binary metadata stubs (indexed as metadata only, no embedding)
    ".shp", ".shx", ".gdbtable", ".gdbtablx", ".gph", ".gpkg",
    ".atx", ".gdbindexes", ".sbn", ".sbx",
    ".spx", ".horizon", ".freelist", ".cfs",
    # Images (metadata stub)
    ".jpg", ".png", ".gif",
}

SKIP_DIRS = {
    ".git", ".venv", "venv", "node_modules", "__pycache__",
    # EDITH source code — don't pollute research index with app code
    "EDITH_M4", "edith_safe_chat", "server", "renderer",
    # Fine-tuning / training data
    "Fine_Tuning", "fine_tuning", "training_data",
    # Backups (47+ snapshots on Edith Bolt)
    "backups",
    # GIS geodatabases (binary, hundreds of internal files)
    # Matched by suffix below, but skip entire .gdb dirs for speed
}

# Directories ending in .gdb are ArcGIS geodatabases — skip entirely
_GDB_SUFFIX = ".gdb"

# Built-in path exclusions (checked against relative path)
EXCLUDE_PATH_PATTERNS = [
    "EDITH_M4/", "edith_safe_chat/",
    "Forge/Fine_Tuning", "Fine_Tuning/",
    "backups/",
]

# §1.8: Lock file to prevent concurrent indexing runs
_LOCK_FILE = None

def _acquire_lock(app_state: Path) -> bool:
    """Acquire a lock file to prevent concurrent indexing runs."""
    global _LOCK_FILE
    lock_path = app_state / "edith_index.lock"
    try:
        if lock_path.exists():
            # Check if the lock is stale (older than 2 hours)
            age = _time.time() - lock_path.stat().st_mtime
            if age < 7200:
                pid = lock_path.read_text().strip()
                print(f"  [lock] Another indexing run is active (PID {pid}, {int(age)}s ago)")
                return False
            print(f"  [lock] Stale lock found ({int(age)}s old), overriding")
        lock_path.write_text(str(os.getpid()))
        _LOCK_FILE = lock_path
        return True
    except Exception as e:
        print(f"  [lock] Warning: {e}")
        return True  # Don't block on lock errors

def _release_lock():
    """Release the indexing lock file."""
    global _LOCK_FILE
    if _LOCK_FILE and _LOCK_FILE.exists():
        try:
            _LOCK_FILE.unlink()
        except Exception:
            pass
        _LOCK_FILE = None

# §2.3: Exclude patterns from EDITH_EXCLUDE_PATTERNS
def _load_exclude_patterns() -> list:
    """Load glob-based exclude patterns from env."""
    import fnmatch
    raw = os.environ.get("EDITH_EXCLUDE_PATTERNS", "")
    if not raw.strip():
        return []
    return [p.strip() for p in raw.split(",") if p.strip()]

def _should_exclude(filename: str, rel_path: str, patterns: list) -> bool:
    """Check if a file matches any exclude pattern."""
    import fnmatch
    for pat in patterns:
        if fnmatch.fnmatch(filename, pat) or fnmatch.fnmatch(rel_path, pat):
            return True
    return False


def setup():
    """Phase 1: Load config, init ChromaDB, init embedder, load manifests.
    Returns a config dict used by scan_and_index() and finalize()."""
    load_env()
    start_time = _time.time()

    docs_root = Path(os.environ.get("EDITH_DATA_ROOT", "")).expanduser().resolve()
    pdf_only = os.environ.get("EDITH_PDF_ONLY", "false").strip().lower() == "true"
    valid_ext = {".pdf"} if pdf_only else set(VALID_EXT)
    default_root = Path(__file__).parent.parent
    app_state = Path(os.environ.get("EDITH_APP_DATA_DIR", str(default_root))).expanduser().resolve()
    chroma_dir = Path(os.environ.get("EDITH_CHROMA_DIR", str(app_state / "chroma"))).expanduser().resolve()
    collection_name = os.environ.get("EDITH_CHROMA_COLLECTION", "edith_docs")
    embed_model = os.environ.get("EDITH_EMBED_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
    report_path = app_state / "edith_index_report.csv"

    chunk_size = max(200, int(os.environ.get("EDITH_CHROMA_CHUNK_CHARS", "1800")))
    overlap = max(0, int(os.environ.get("EDITH_CHROMA_CHUNK_OVERLAP_CHARS", "250")))
    if overlap >= chunk_size:
        overlap = max(0, chunk_size // 6)
    prune_stale = os.environ.get("EDITH_CHROMA_PRUNE_STALE", "true").lower() == "true"
    ocr_on_ingest = os.environ.get("EDITH_OCR_ON_INGEST", "false").lower() == "true"
    content_aware_chunking = os.environ.get("EDITH_CONTENT_AWARE_CHUNKS", "true").lower() == "true"
    dedup_enabled = os.environ.get("EDITH_DEDUP_ENABLED", "true").lower() == "true"
    try:
        ocr_max_pages = int(os.environ.get("EDITH_OCR_MAX_PAGES", "50"))
    except Exception:
        ocr_max_pages = 50

    # ── Safety: prevent indexing from broad root dirs ──
    if os.environ.get("EDITH_ALLOW_WIDE_ROOT", "false").strip().lower() != "true":
        home = Path.home().resolve()
        blocked_exact = {home, home / "Desktop", home / "Documents", home / "Downloads", home / "Library"}
        blocked_names = {"my drive", "google drive", "icloud drive", "onedrive", "cloudstorage"}
        if docs_root in blocked_exact or docs_root.name.strip().lower() in blocked_names:
            raise SystemExit(
                "Unsafe EDITH_DATA_ROOT: choose a dedicated folder (example: .../My Drive/Edith_Index), "
                "not a broad root."
            )

    if not docs_root.exists() or not docs_root.is_dir():
        raise SystemExit(f"EDITH_DATA_ROOT not found: {docs_root}")

    # §1.1: Config validation — fail fast if no indexable files
    quick_check = False
    for _r, _d, _f in os.walk(docs_root):
        _d[:] = [d for d in _d if d not in SKIP_DIRS]
        for fn in _f:
            if not fn.startswith(".") and Path(fn).suffix.lower() in valid_ext:
                quick_check = True
                break
        if quick_check:
            break
    if not quick_check:
        raise SystemExit(f"No indexable files found in {docs_root}. Check EDITH_DATA_ROOT.")

    app_state.mkdir(parents=True, exist_ok=True)
    chroma_dir.mkdir(parents=True, exist_ok=True)

    # §1.8: Acquire lock
    if not _acquire_lock(app_state):
        raise SystemExit("Another indexing run is active. Delete edith_index.lock to force.")

    # §1.2: Dry-run mode
    dry_run = os.environ.get("EDITH_DRY_RUN", "false").lower() == "true"
    if "--dry-run" in (os.environ.get("EDITH_INDEX_ARGS", "")).split():
        dry_run = True

    # ── ChromaDB collections ──
    client = chromadb.PersistentClient(path=str(chroma_dir))
    collection = client.get_or_create_collection(name=collection_name, metadata={"hnsw:space": "cosine"})
    section_collection = client.get_or_create_collection(name=f"{collection_name}_sections", metadata={"hnsw:space": "cosine"})
    # Fix #6: Metadata-only collection for GIS stubs (no embedding waste)
    metadata_collection = client.get_or_create_collection(name=f"{collection_name}_metadata", metadata={"hnsw:space": "cosine"})

    embedder = SentenceTransformer(embed_model)

    # ── Load unified manifest (Fix #4: single manifest instead of two) ──
    manifest_path = app_state / "edith_index_manifest.json"
    file_manifest: dict = {}  # {rel_path: {mtime, sha256}}
    known_hashes: set = set()
    if manifest_path.exists():
        try:
            data = json.loads(manifest_path.read_text(encoding="utf-8"))
            file_manifest = data.get("files", {})
            known_hashes = set(data.get("known_hashes", []))
            print(f"  [manifest] loaded {len(file_manifest)} files, {len(known_hashes)} hashes")
        except Exception as e:
            print(f"  [manifest] failed to load: {e}")
    try:
        # If this collection is empty (new collection name), do not trust manifest
        # hashes from previous collections/runs.
        if collection.count() == 0 and known_hashes:
            print("  [manifest] collection is empty — forcing full re-index for this collection")
            known_hashes = set()
    except Exception:
        pass
    if not known_hashes:
        # Fallback: scan Chroma (first run or corrupted manifest)
        print("  [cache] no manifest hashes — scanning Chroma for existing signatures...")
        try:
            count = collection.count()
            PAGE = 10000
            for offset in range(0, count, PAGE):
                res = collection.get(include=["metadatas"], limit=PAGE, offset=offset)
                for m in (res.get("metadatas") or []):
                    s = m.get("sha256")
                    if s:
                        known_hashes.add(s)
            print(f"  [cache] {len(known_hashes)} unique files found in Chroma.")
        except Exception as e:
            print(f"  [cache] warning: could not scan Chroma ({e})")
    print(f"  [cache] {len(known_hashes)} unique files already indexed.")

    # ── Embedding strategy: Gemini API > local SentenceTransformer ──
    use_gemini_embed = os.environ.get("EDITH_USE_GEMINI_EMBED", "true").lower() == "true"
    gemini_embed_model = os.environ.get("EDITH_GEMINI_EMBED_MODEL", "gemini-embedding-001")
    embed_dim = None
    if use_gemini_embed:
        test = _gemini_embed_batch(["test"], model=gemini_embed_model)
        if test and len(test) == 1 and len(test[0]) > 0:
            embed_dim = len(test[0])
            print(f"  [embed] Using Gemini API ({gemini_embed_model}, {embed_dim}-dim)")
        else:
            print(f"  [embed] Gemini API unavailable, falling back to local {embed_model}")
            use_gemini_embed = False
    if not use_gemini_embed:
        print(f"  [embed] Using local model: {embed_model}")

    # ── Dimension safety ──
    try:
        existing_count = collection.count()
        if existing_count > 0 and embed_dim:
            sample = collection.get(limit=1, include=["embeddings"])
            existing_embeds = (sample.get("embeddings") or [[]])[0]
            existing_dim = len(existing_embeds) if existing_embeds else 0
            if existing_dim > 0 and existing_dim != embed_dim:
                print(f"  [DIMENSION MISMATCH] Existing: {existing_dim}-dim, new: {embed_dim}-dim")
                collection_name = f"{collection_name}_v2"
                collection = client.get_or_create_collection(name=collection_name, metadata={"hnsw:space": "cosine"})
                section_collection = client.get_or_create_collection(name=f"{collection_name}_sections", metadata={"hnsw:space": "cosine"})
                known_hashes = set()
                print(f"  [FRESH] Will re-index all files into {collection_name}")
    except Exception as e:
        print(f"  [dim-check] warning: {e}")

    # ── Improvement 1.3/1.5/1.8/1.9: Checkpoint, Dedup, Progress, Validation ──
    checkpoint = None
    progress = None
    validator = None
    dedup = None
    embed_version = {}
    if _IMPROVEMENTS_AVAILABLE:
        checkpoint = IndexCheckpoint(app_state / "edith_index_checkpoint.json")
        progress = IndexProgressBroadcaster(app_state / "edith_index_progress.json")
        validator = IngestValidator()
        if dedup_enabled:
            dedup = DuplicateDetector(threshold=3)
        embed_version = get_embed_version_info(
            use_gemini_embed, gemini_embed_model, embed_model
        )
        if checkpoint.resumed:
            print(f"  [checkpoint] Resuming from: {checkpoint.stats}")
    else:
        print("  [info] indexing_improvements module not available — basic mode")

    # §2.3: Load exclude patterns
    exclude_patterns = _load_exclude_patterns()
    if exclude_patterns:
        print(f"  [exclude] Patterns: {', '.join(exclude_patterns)}")
    if pdf_only:
        print("  [mode] PDF-only indexing enabled")

    # §AI CLASSIFICATION: Load classification cache
    _load_ai_classify_cache(app_state)

    return {
        "docs_root": docs_root,
        "app_state": app_state,
        "chroma_dir": chroma_dir,
        "collection_name": collection_name,
        "report_path": report_path,
        "chunk_size": chunk_size,
        "overlap": overlap,
        "prune_stale": prune_stale,
        "ocr_on_ingest": ocr_on_ingest,
        "ocr_max_pages": ocr_max_pages,
        "content_aware_chunking": content_aware_chunking,
        "client": client,
        "collection": collection,
        "section_collection": section_collection,
        "metadata_collection": metadata_collection,
        "embedder": embedder,
        "manifest_path": manifest_path,
        "file_manifest": file_manifest,
        "known_hashes": known_hashes,
        "use_gemini_embed": use_gemini_embed,
        "gemini_embed_model": gemini_embed_model,
        "start_time": start_time,
        "dry_run": dry_run,
        "exclude_patterns": exclude_patterns,
        "valid_ext": valid_ext,
        "pdf_only": pdf_only,
        # Improvements
        "checkpoint": checkpoint,
        "progress": progress,
        "validator": validator,
        "dedup": dedup,
        "embed_version": embed_version,
    }


def _get_sha256_incremental(path: Path, file_manifest: dict) -> str:
    """Return SHA256 for a file, using manifest cache when mtime hasn't changed."""
    key = str(path)
    try:
        current_mtime = int(path.stat().st_mtime)
    except Exception:
        h = sha256_file(path)
        return h if h else ""
    cached = file_manifest.get(key)
    if cached and cached.get("mtime") == current_mtime and cached.get("sha256"):
        return cached["sha256"]
    h = sha256_file(path)
    if h is None:
        return ""  # File unreadable (timeout, permission, etc.) — skip it
    file_manifest[key] = {"mtime": current_mtime, "sha256": h}
    return h


def _save_unified_manifest(manifest_path: Path, file_manifest: dict, known_hashes: set,
                           docs_root: str, chroma_dir: str):
    """Save unified manifest: file-level mtime/sha256 + known hashes set. (Fix #4)"""
    try:
        data = {
            "files": file_manifest,
            "known_hashes": sorted(known_hashes),
            "count": len(known_hashes),
            "file_count": len(file_manifest),
            "docs_root": docs_root,
            "chroma_dir": chroma_dir,
            "saved_at": _time.strftime("%Y-%m-%dT%H:%M:%S"),
        }
        manifest_path.write_text(json.dumps(data, indent=1), encoding="utf-8")
        print(f"  [manifest] saved {len(file_manifest)} files, {len(known_hashes)} hashes")
    except Exception as e:
        print(f"  [manifest] warning: could not save ({e})")


def scan_and_index(cfg: dict):
    """Phase 2: Walk files, extract text, chunk, embed, upsert to ChromaDB."""
    docs_root = cfg["docs_root"]
    collection = cfg["collection"]
    section_collection = cfg["section_collection"]
    metadata_collection = cfg["metadata_collection"]
    embedder = cfg["embedder"]
    file_manifest = cfg["file_manifest"]
    known_hashes = cfg["known_hashes"]
    chunk_size = cfg["chunk_size"]
    overlap = cfg["overlap"]
    ocr_on_ingest = cfg["ocr_on_ingest"]
    ocr_max_pages = cfg["ocr_max_pages"]
    use_gemini_embed = cfg["use_gemini_embed"]
    gemini_embed_model = cfg["gemini_embed_model"]
    manifest_path = cfg["manifest_path"]
    content_aware_chunking = cfg.get("content_aware_chunking", False)
    checkpoint = cfg.get("checkpoint")
    progress = cfg.get("progress")
    validator = cfg.get("validator")
    dedup = cfg.get("dedup")
    embed_version = cfg.get("embed_version", {})
    valid_ext = set(cfg.get("valid_ext") or VALID_EXT)
    dry_run = cfg.get("dry_run", False)
    exclude_patterns = cfg.get("exclude_patterns", [])

    if dry_run:
        print("\n  *** DRY-RUN MODE: No writes to ChromaDB ***\n")

    # Start progress broadcast
    if progress:
        progress.start(total_estimate=0)

    def _embed_texts(text_list, task_type="RETRIEVAL_DOCUMENT"):
        """Embed texts using Gemini API with retry, or local model.
        Uses parallel embedding (improvement 1.2) when available."""
        if use_gemini_embed:
            # Try parallel embedding first (improvement 1.2)
            if _IMPROVEMENTS_AVAILABLE and len(text_list) > 20:
                vecs = parallel_gemini_embed(
                    text_list, model=gemini_embed_model,
                    task_type=task_type, max_concurrent=3,
                )
                if vecs and len(vecs) == len(text_list):
                    return vecs
            # Fallback to serial embedding with retry
            for attempt in range(4):
                vecs = _gemini_embed_batch(text_list, model=gemini_embed_model, task_type=task_type)
                if vecs and len(vecs) == len(text_list):
                    return vecs
                if attempt < 3:
                    wait = 2 ** (attempt + 1)
                    print(f"  [embed] Gemini attempt {attempt+1} failed, retrying in {wait}s...")
                    _time.sleep(wait)
            existing_count = collection.count()
            if existing_count > 0:
                print(f"  [embed] WARNING: Gemini failed after 4 attempts. Skipping batch ({len(text_list)} texts) to avoid dimension mismatch.")
                return None
            print(f"  [embed] Gemini unavailable, falling back to local model (collection is empty)")
        return embedder.encode(text_list, normalize_embeddings=True, show_progress_bar=True).tolist()

    # Batch buffers
    ids, texts, metas = [], [], []
    s_ids, s_texts, s_metas = [], [], []
    m_ids, m_docs, m_metas = [], [], []

    active_hashes = set()
    total_chunks = 0
    total_sections = 0
    total_metadata_stubs = 0
    total_errors = 0
    total_duplicates = 0

    def flush_batch():
        nonlocal ids, texts, metas, s_ids, s_texts, s_metas, total_chunks, total_sections
        nonlocal m_ids, m_docs, m_metas, total_metadata_stubs
        # §1.2: Skip writes in dry-run mode
        if dry_run:
            total_chunks += len(texts)
            total_sections += len(s_texts)
            total_metadata_stubs += len(m_ids)
            print(f"  [dry-run] Would flush {len(texts)} chunks, {len(s_texts)} sections, {len(m_ids)} stubs")
            ids, texts, metas = [], [], []
            s_ids, s_texts, s_metas = [], [], []
            m_ids, m_docs, m_metas = [], [], []
            return
        if texts:
            vectors = _embed_texts(texts)
            if vectors is None:
                print(f"  [batch] SKIPPED {len(texts)} chunks (Gemini unavailable)")
            else:
                for i in range(0, len(ids), 5000):
                    collection.upsert(
                        ids=ids[i:i+5000],
                        documents=texts[i:i+5000],
                        embeddings=vectors[i:i+5000],
                        metadatas=metas[i:i+5000],
                    )
                total_chunks += len(texts)
                print(f"  [batch] flushed {len(texts)} chunks (total: {total_chunks})")
            ids, texts, metas = [], [], []

        if s_texts:
            s_vectors = _embed_texts(s_texts)
            if s_vectors is None:
                print(f"  [batch] SKIPPED {len(s_texts)} sections (Gemini unavailable)")
            else:
                for i in range(0, len(s_ids), 5000):
                    section_collection.upsert(
                        ids=s_ids[i:i+5000],
                        documents=s_texts[i:i+5000],
                        embeddings=s_vectors[i:i+5000],
                        metadatas=s_metas[i:i+5000],
                    )
                total_sections += len(s_texts)
                print(f"  [batch] flushed {len(s_texts)} sections (total: {total_sections})")
            s_ids, s_texts, s_metas = [], [], []

        # Fix #6: GIS metadata stubs — upsert WITHOUT embeddings
        if m_ids:
            for i in range(0, len(m_ids), 5000):
                metadata_collection.upsert(
                    ids=m_ids[i:i+5000],
                    documents=m_docs[i:i+5000],
                    metadatas=m_metas[i:i+5000],
                )
            total_metadata_stubs += len(m_ids)
            print(f"  [batch] flushed {len(m_ids)} metadata stubs (total: {total_metadata_stubs})")
            m_ids, m_docs, m_metas = [], [], []

        # Fix #5: Save manifest inside flush for crash resilience
        all_hashes = known_hashes | active_hashes
        _save_unified_manifest(manifest_path, file_manifest, all_hashes,
                               str(docs_root), str(cfg["chroma_dir"]))

    rows = []
    vault_manifest = load_vault_file_manifest(docs_root)
    BATCH_FILES = 50
    files_seen = 0
    files_since_flush = 0
    current_folder = ""

    # §FIX: Single-file mode — skip the full walk, only process one file
    single_file = cfg.get("single_file")
    if single_file:
        _sf_path = Path(single_file)
        if _sf_path.exists():
            _sf_files_iter = [(_sf_path.parent, [], [_sf_path.name])]
            print(f"\n  [single-file] Indexing only: {_sf_path.name}")
        else:
            print(f"  [single-file] File not found, falling back to full scan")
            _sf_files_iter = None
            single_file = None
    else:
        _sf_files_iter = None

    _walk_iter = _sf_files_iter if _sf_files_iter else os.walk(docs_root, followlinks=True)

    for root, dirs, files in _walk_iter:
        if not single_file:
            dirs[:] = [
                d for d in dirs
                if d not in SKIP_DIRS and not d.lower().endswith(_GDB_SUFFIX)
            ]

        # §FIX Robustness 2: Handle single-file outside docs_root gracefully
        try:
            rel_folder = str(Path(root).relative_to(docs_root))
        except ValueError:
            rel_folder = str(Path(root))
        if rel_folder != current_folder:
            print(f"\n--- SCANNING FOLDER: {rel_folder or 'root'} ---")
            current_folder = rel_folder

        for fn in files:
            if fn.startswith("."):
                continue
            path = Path(root) / fn

            # §FIX: Symlink safety — skip broken symlinks (bus error on dangling links)
            if path.is_symlink() and not path.exists():
                print(f"  [skip] broken symlink: {fn}")
                continue

            ext = path.suffix.lower()
            if ext not in valid_ext:
                continue

            # §2.3: Apply exclude patterns
            try:
                rel = str(path.relative_to(docs_root))
            except ValueError:
                rel = str(path)

            if exclude_patterns and _should_exclude(fn, rel, exclude_patterns):
                continue

            # Built-in path exclusions
            if any(excl in rel for excl in EXCLUDE_PATH_PATTERNS):
                continue

            files_seen += 1
            if files_seen % 50 == 0:
                print(f"  [scan] Processed {files_seen} files...")

            # §FIX: Catch unreadable files (permission errors, etc.)
            try:
                f_sha = _get_sha256_incremental(path, file_manifest)
            except (OSError, IOError) as e:
                print(f"  [skip] unreadable: {fn} ({e})")
                continue

            # Improvement 1.3: Skip if already checkpointed in this run
            if checkpoint and checkpoint.is_processed(f_sha):
                active_hashes.add(f_sha)
                continue

            # Skip already-indexed files
            if f_sha in known_hashes:
                if f_sha not in active_hashes:
                    active_hashes.add(f_sha)
                    print(f"  [skip] {path.name} (already indexed)")
                continue
            if f_sha in active_hashes:
                continue
            active_hashes.add(f_sha)

            print(f"  [index] Extracting & Vectorizing: {path.name} ({round(os.path.getsize(path)/1024, 1)} KB)")

            rel = str(path.relative_to(docs_root))
            project = infer_project(rel)
            academic_topic = infer_academic_topic(rel)

            # Tiering
            tier = "inbox"
            rel_lower = rel.lower()
            if "canon" in rel_lower or "/canon" in rel_lower:
                tier = "canon"
            elif "project" in rel_lower:
                tier = "projects"

            tag = infer_tag(path.name)
            title_guess, author_guess, year_guess = infer_citation(path.name)
            citation_source = "filename"
            doc_type = infer_doc_type(rel, path.suffix)
            version_stage = infer_version_stage(path.name)
            if ext == ".pdf":
                try:
                    pdf_title, pdf_author, pdf_year, pdf_source = infer_citation_from_pdf(path)
                    if pdf_title: title_guess = pdf_title
                    if pdf_author: author_guess = pdf_author
                    if pdf_year: year_guess = pdf_year
                    if pdf_source != "filename": citation_source = pdf_source
                except Exception:
                    pass
            doc_family = infer_doc_family(title_guess, author_guess, year_guess, rel)
            vault_meta = vault_manifest.get(rel) or {}

            # Fix #6: GIS binary files → metadata stub only (no embedding)
            if ext in GIS_BINARY_EXTS:
                stub_blocks = extract_gis_metadata_stub(path)
                for sb in stub_blocks:
                    stub_text = sb.get("text", "")
                    m_ids.append(f"{f_sha}:meta:0")
                    m_docs.append(stub_text)
                    m_metas.append({
                        "path": rel, "file_name": path.name, "sha256": f_sha,
                        "project": project, "tier": tier, "doc_type": doc_type,
                        "file_type": sb.get("section_heading", "Binary file"),
                    })
                files_since_flush += 1
                if files_since_flush >= BATCH_FILES:
                    flush_batch()
                    files_since_flush = 0
                continue

            # Regular text extraction (with enhanced OCR - improvement 1.1)
            try:
                blocks = extract_text_blocks(path, ocr_on_ingest=ocr_on_ingest, ocr_max_pages=ocr_max_pages)
                # Improvement 1.1: If no text extracted and OCR available, try enhanced OCR
                if not blocks and _IMPROVEMENTS_AVAILABLE and ext == ".pdf":
                    blocks = enhanced_ocr_extract(path, ocr_max_pages=ocr_max_pages)
            except Exception as exc:
                total_errors += 1
                if validator:
                    validator.record_failure(rel, "extract_error", str(exc))
                print(f"SKIP (extract error): {rel}: {exc}")
                continue
            if not blocks:
                print(f"  [WARN] No extractable text from: {rel} (ext={ext})")
                if validator:
                    validator.record_failure(rel, "empty_text", "No extractable text")
                continue

            # Improvement 1.5: Duplicate detection
            all_text = " ".join([b.get("text") or "" for b in blocks])
            if dedup:
                dup_of = dedup.check(f_sha, all_text)
                if dup_of:
                    total_duplicates += 1
                    if validator:
                        validator.record_warning(rel, f"Near-duplicate of {dup_of}")
                    print(f"  [dedup] {path.name} is near-duplicate (hamming ≤ 3), indexing anyway")

            # §AI CLASSIFICATION: Use Gemini to auto-classify the paper
            ai_class = ai_classify_paper(
                text_sample=all_text[:2500],
                sha256=f_sha,
                filename=path.name,
                existing_topic=academic_topic,
            )
            # Override folder-based topic if it's vague and AI found something better
            if ai_class.get("topic") and academic_topic.lower() in {"inbox", "general", "", "uncategorized", "root"}:
                academic_topic = ai_class["topic"]
                print(f"  [ai-classify] {path.name} → topic: {academic_topic}")
            # Always store method, country, theory as additional metadata
            ai_method = ai_class.get("method", "")
            ai_country = ai_class.get("country", "")
            ai_theory = ai_class.get("theory", "")
            ai_doc_type = ai_class.get("doc_type_refined", "")
            if ai_method or ai_country or ai_theory:
                print(f"  [ai-classify] {path.name} → method: {ai_method}, country: {ai_country}, theory: {ai_theory}")

            # Improvement 1.10: Detect language
            doc_language = "en"
            if _IMPROVEMENTS_AVAILABLE:
                doc_language = detect_language(all_text)

            per_file_chunk = 0
            file_ocr_used = False
            total_words = len(all_text.split())
            file_sections = {}

            # Improvement 1.4: Content-aware chunk sizing
            if content_aware_chunking and _IMPROVEMENTS_AVAILABLE:
                chunk_size_eff, overlap_eff = content_aware_chunk_params(
                    doc_type, ext, total_words
                )
            else:
                chunk_size_eff, overlap_eff = chunk_size, overlap

            for block in blocks:
                page = block.get("page")
                section_heading = clean_text(block.get("section_heading") or "")
                ocr_used = bool(block.get("ocr_used"))
                if ocr_used:
                    file_ocr_used = True

                block_text = block.get("text") or ""
                if section_heading:
                    file_sections[section_heading] = file_sections.get(section_heading, "") + "\n" + block_text

                block_chunks = chunk_text(block_text, chunk_size_eff, overlap_eff)

                # Improvement 1.6: Extract structured tables
                if _IMPROVEMENTS_AVAILABLE:
                    tables = extract_table_from_text(block_text)
                    if tables:
                        table_texts = tables_to_searchable_text(tables)
                        for tt in table_texts:
                            block_chunks.append(tt)

                for ch in block_chunks:
                    figure_table_markers = extract_figure_table_markers(ch)
                    equation_markers = extract_equation_markers(ch)
                    doc_id = f"{f_sha}:{per_file_chunk}"
                    per_file_chunk += 1
                    ids.append(doc_id)
                    # Contextual header for better embedding quality
                    ctx_parts = []
                    if title_guess: ctx_parts.append(title_guess)
                    if section_heading: ctx_parts.append(section_heading)
                    ctx_header = " | ".join(ctx_parts)
                    enriched_chunk = f"[{ctx_header}] {ch}" if ctx_header else ch
                    texts.append(enriched_chunk)
                    # §FIX: Sanitize metadata — ChromaDB rejects None values
                    _page_val = int(page) if page is not None else 0
                    metas.append({
                        "path": rel,
                        "rel_path": rel,
                        "file_name": path.name,
                        "project": project or "",
                        "academic_topic": academic_topic or "",
                        "tag": tag or "",
                        "tier": tier or "",
                        "sha256": f_sha,
                        "total_words": total_words,
                        "chunk": per_file_chunk - 1,
                        "page": _page_val,
                        "section_heading": section_heading or "",
                        "doc_type": ai_doc_type or doc_type or "",
                        "version_stage": version_stage or "",
                        "author": author_guess or "",
                        "year": year_guess or "",
                        "title": title_guess or "",
                        "citation_source": citation_source or "",
                        "figure_table_markers": ",".join(figure_table_markers),
                        "equation_markers": ",".join(equation_markers),
                        # §AI CLASSIFICATION metadata
                        "method": ai_method or "",
                        "country": ai_country or "",
                        "theory": ai_theory or "",
                        "ocr_used": ocr_used,
                        "doc_family": doc_family or "",
                        "language": doc_language or "en",
                        "embed_model": embed_version.get("embed_model", ""),
                        "embed_dim": embed_version.get("embed_dim", 0),  # Improvement 1.7
                        "vault_export_id": join_manifest_field(vault_meta.get("export_id")),
                        "vault_export_date": join_manifest_field(vault_meta.get("export_date")),
                        "vault_custodian": join_manifest_field(vault_meta.get("custodian")),
                        "vault_matter_name": join_manifest_field(vault_meta.get("matter_name")),
                        "quality_score": _chunk_quality_score(ch),  # §4.0: info density
                    })

            # Section summaries
            for heading, s_text in file_sections.items():
                if len(s_text) < 50:
                    continue
                s_ids.append(f"{f_sha}:sec:{hash(heading)}")
                s_texts.append(s_text[:8000])
                s_metas.append({
                    "path": rel, "sha256": f_sha, "tier": tier,
                    "section_heading": heading, "doc_type": doc_type,
                    "title": title_guess, "author": author_guess, "year": year_guess,
                })

            rows.append({
                "file_name": path.name, "rel_path": rel, "project": project,
                "tag": tag, "title_guess": title_guess, "author_guess": author_guess,
                "year_guess": year_guess, "doc_type": doc_type,
                "version_stage": version_stage, "doc_family": doc_family,
                "ocr_used": str(bool(file_ocr_used)).lower(),
                "citation_source": citation_source,
                "language": doc_language,
                "vault_export_id": join_manifest_field(vault_meta.get("vault_export_id")),
                "vault_export_date": join_manifest_field(vault_meta.get("vault_export_date")),
                "vault_custodian": join_manifest_field(vault_meta.get("vault_custodian")),
                "vault_matter_name": join_manifest_field(vault_meta.get("vault_matter_name")),
            })

            # Improvement 1.3: Checkpoint after each file
            if checkpoint:
                checkpoint.mark_processed(f_sha)
            # Improvement 1.9: Record success
            if validator:
                validator.record_success(rel, per_file_chunk, file_ocr_used)

            files_since_flush += 1
            if files_since_flush >= BATCH_FILES:
                flush_batch()
                files_since_flush = 0
                # Improvement 1.3: Save checkpoint on flush
                if checkpoint:
                    checkpoint.update_progress(current_folder, files_seen, total_chunks)
                    checkpoint.save()
                # Improvement 1.8: Broadcast progress
                if progress:
                    progress.update(
                        file_name=path.name, folder=current_folder,
                        processed=files_seen, chunks=total_chunks,
                        errors=total_errors, phase="indexing",
                    )
                print(f"  [progress] {files_seen} files scanned so far")

    # Final flush
    flush_batch()

    # §AI CLASSIFICATION: Save cache to disk
    _save_ai_classify_cache()

    # Improvement 1.3: Clear checkpoint on success
    if checkpoint:
        checkpoint.clear()
    # Improvement 1.9: Generate validation report
    if validator:
        report_dir = cfg["app_state"]
        validator.generate_report(report_dir / "edith_ingest_validation.json")
        stats = validator.stats
        print(f"  [validation] {stats['successful']} ok, {stats['failed']} failed, "
              f"{stats['ocr_used']} OCR, {stats['warnings']} warnings")
    # Improvement 1.5: Duplicate report
    if dedup and dedup.duplicates:
        print(f"  [dedup] Found {len(dedup.duplicates)} near-duplicate files")
    # Improvement 1.8: Broadcast completion
    if progress:
        progress.complete({
            "files_processed": files_seen, "chunks_added": total_chunks,
            "duplicates": total_duplicates, "errors": total_errors,
        })

    return {
        "files_seen": files_seen,
        "rows": rows,
        "active_hashes": active_hashes,
        "total_chunks": total_chunks,
        "total_sections": total_sections,
        "total_metadata_stubs": total_metadata_stubs,
        "total_duplicates": total_duplicates,
        "total_errors": total_errors,
        "validation_report": validator.stats if validator else {},
        "duplicate_report": dedup.report if dedup else [],
    }


def finalize(cfg: dict, results: dict):
    """Phase 3: Prune stale entries, write CSV report, print summary."""
    collection = cfg["collection"]
    section_collection = cfg["section_collection"]
    report_path = cfg["report_path"]
    prune_stale = cfg["prune_stale"]
    active_hashes = results["active_hashes"]
    rows = results["rows"]
    total_chunks = results["total_chunks"]
    total_sections = results["total_sections"]
    total_metadata_stubs = results["total_metadata_stubs"]
    files_seen = results["files_seen"]

    if total_chunks == 0 and total_metadata_stubs == 0:
        print(json.dumps({"ok": True, "message": "No indexable content found.", "files_seen": files_seen}, indent=2))
        return

    # Fix #3: Paginated stale pruning (avoids OOM with 260k+ entries)
    stale_deleted = 0
    if prune_stale:
        try:
            count = collection.count()
            PAGE = 10000
            stale_ids = []
            stale_filenames = set()  # §FIX I1: Collect filenames in first pass
            for offset in range(0, count, PAGE):
                batch = collection.get(include=["metadatas"], limit=PAGE, offset=offset)
                batch_ids = flatten_if_nested(batch.get("ids") or [])
                batch_meta = flatten_if_nested(batch.get("metadatas") or [])
                for rid, meta in zip(batch_ids, batch_meta):
                    sha = str((meta or {}).get("sha256") or "").strip()
                    if sha and sha not in active_hashes:
                        stale_ids.append(rid)
                        stale_filenames.add((meta or {}).get("file_name", "unknown"))
            if stale_ids:
                # §9.2: Pruning dry-run — show what would be deleted
                print(f"  [prune] Found {len(stale_ids)} stale chunk(s) to remove")
                for name in sorted(stale_filenames)[:10]:
                    print(f"    • {name}")
                if len(stale_filenames) > 10:
                    print(f"    … and {len(stale_filenames) - 10} more files")

                if cfg.get("dry_run"):
                    print(f"  [prune] DRY-RUN: would delete {len(stale_ids)} chunks")
                else:
                    for i in range(0, len(stale_ids), 1000):
                        collection.delete(ids=stale_ids[i:i + 1000])
                    stale_deleted = len(stale_ids)
                    print(f"  [prune] Deleted {stale_deleted} stale chunks")
        except Exception as e:
            print(f"  [prune] warning: {e}")
            stale_deleted = 0

    # Write CSV report
    with report_path.open("w", newline="") as f:
        writer = csv.DictWriter(
            f,
            fieldnames=[
                "file_name", "rel_path", "project", "tag",
                "title_guess", "author_guess", "year_guess",
                "doc_type", "version_stage", "doc_family",
                "ocr_used", "citation_source", "language",
                "vault_export_id", "vault_export_date",
                "vault_custodian", "vault_matter_name",
            ],
            extrasaction="ignore",  # Skip unknown fields instead of crashing
        )
        writer.writeheader()
        writer.writerows(rows)

    # ── Post-index hook: feed chunks through ScholarlyRepositories + CitadelGraph ──
    scholarly_stats = {}
    graph_stats = {}
    try:
        from server.vault_config import VAULT_ROOT
        from server.scholarly_repositories import ScholarlyRepositories
        repos = ScholarlyRepositories(store_dir=str(VAULT_ROOT / "Corpus" / "Vault" / "scholarly"))
        # Sample chunks from newly-indexed docs
        sample_size = min(500, total_chunks)
        if sample_size > 0:
            try:
                batch = collection.get(limit=sample_size, include=["documents", "metadatas"])
                fed = 0
                for i, doc in enumerate(batch.get("documents") or []):
                    if not doc or len(doc) < 100:
                        continue
                    meta = (batch.get("metadatas") or [{}])[i] or {}
                    paper_id = (batch.get("ids") or [f"chunk_{i}"])[i]
                    repos.process_chunk(doc, paper_id, {
                        "title_guess": meta.get("title_guess", ""),
                        "author_guess": meta.get("author_guess", ""),
                        "year_guess": str(meta.get("year_guess", "")),
                    })
                    fed += 1
                scholarly_stats = repos.stats()
                print(f"  [scholarly] Fed {fed} chunks → {scholarly_stats['datasets']} datasets, "
                      f"{scholarly_stats['methods']} methods, {scholarly_stats['countries']} countries")
            except Exception as e:
                print(f"  [scholarly] warning: {e}")
    except ImportError:
        print("  [scholarly] ScholarlyRepositories not available")
    except Exception as e:
        print(f"  [scholarly] warning: {e}")

    try:
        from server.graph_vector_engine import GraphVectorEngine
        engine = GraphVectorEngine()
        # Feed the CSV rows as summary documents for entity extraction
        graph_fed = 0
        for row in rows[:100]:  # Cap at 100 docs for performance
            title = row.get("title_guess", "")
            author = row.get("author_guess", "")
            if title:
                text = f"{title} by {author}" if author else title
                try:
                    import tempfile
                    with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
                        f.write(text)
                        tmp = f.name
                    engine.ingest_document(tmp, source_label=row.get("file_name", ""))
                    os.unlink(tmp)
                    graph_fed += 1
                except Exception:
                    pass
        graph_stats = engine.get_graph_stats() if graph_fed > 0 else {}
        if graph_fed > 0:
            print(f"  [citadel] Fed {graph_fed} docs → {graph_stats.get('total_entities', 0)} entities, "
                  f"{graph_stats.get('total_relationships', 0)} relationships")
    except ImportError:
        print("  [citadel] GraphVectorEngine not available")
    except Exception as e:
        print(f"  [citadel] warning: {e}")

    # Final summary
    elapsed = _time.time() - cfg["start_time"]
    final_chunk_count = collection.count()
    final_section_count = section_collection.count()

    summary = {
        "ok": True,
        "files_scanned": files_seen,
        "files_indexed": len(rows),
        "total_chunks_in_db": final_chunk_count,
        "total_sections_in_db": final_section_count,
        "chunks_added_this_run": total_chunks,
        "sections_added_this_run": total_sections,
        "metadata_stubs_added": total_metadata_stubs,
        "stale_deleted": stale_deleted,
        "docs_root": str(cfg["docs_root"]),
        "chroma_dir": str(cfg["chroma_dir"]),
        "collection": cfg["collection_name"],
        "report": str(report_path),
        "duration_minutes": round(elapsed / 60, 1),
    }

    print("\n" + "=" * 60)
    print("  CHROMA INDEXING COMPLETE")
    print("=" * 60)
    print(json.dumps(summary, indent=2))
    print("=" * 60)


def main():
    """Orchestrator: setup → scan → finalize."""
    try:
        cfg = setup()

        # §FIX: Support --single-file for auto-reindex after upload
        import sys as _main_sys
        single_file = None
        if "--single-file" in _main_sys.argv:
            idx = _main_sys.argv.index("--single-file")
            if idx + 1 < len(_main_sys.argv):
                single_file = Path(_main_sys.argv[idx + 1])
                if single_file.exists():
                    print(f"  [single-file] Indexing only: {single_file.name}")
                    cfg["single_file"] = single_file
                else:
                    print(f"  [single-file] File not found: {single_file}")
                    return

        results = scan_and_index(cfg)
        finalize(cfg, results)

        # §9.10: Auto-invalidate retrieval caches after successful re-index
        if results.get("total_chunks", 0) > 0:
            try:
                from server.chroma_backend import _query_cache
                _query_cache.invalidate()
                print("  [cache] Query cache invalidated after re-index")
            except ImportError:
                pass  # server module not available in standalone mode
    finally:
        # §1.8: Always release lock
        _release_lock()


if __name__ == "__main__":
    main()
